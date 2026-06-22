from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor

prs = Presentation(r"c:\Users\Admin\Documents\HomeLodge-Doc\docs\Presentation\PSM1_Development.pptx")

def safe_color(run):
    try:
        c = run.font.color
        if c and c.type:
            return str(c.rgb)
    except:
        pass
    return "inherited"

# Print slides 1-16 only
for si, slide in enumerate(prs.slides):
    if si >= 16:
        break
    layout_name = slide.slide_layout.name if slide.slide_layout else "?"
    print(f"\n=== SLIDE {si+1} (layout: {layout_name}) ===")
    for shape in slide.shapes:
        stype = str(shape.shape_type)
        pos = f"({shape.left/914400:.2f}\", {shape.top/914400:.2f}\")"
        size = f"{shape.width/914400:.2f}\"x{shape.height/914400:.2f}\""
        print(f"  [{stype}] '{shape.name}'  pos={pos}  size={size}")
        if shape.has_text_frame:
            for pi, para in enumerate(shape.text_frame.paragraphs):
                txt = para.text.strip()
                if not txt:
                    continue
                font_info = []
                for run in para.runs:
                    font_info.append(f"bold={run.font.bold} sz={run.font.size} col={safe_color(run)}")
                print(f"    [{pi}] '{txt[:120]}' | {'; '.join(font_info[:3])}")
