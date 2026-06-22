"""
HomeLodge PSM I — Presentation v3 (New Version)
Uses PSM1_Development.pptx as the template.
Follows the PSM1 format guidelines exactly: 10 slides.
"""
import sys
sys.stdout.reconfigure(encoding='utf-8')

import copy
import lxml.etree as etree
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml

TEMPLATE = r"c:\Users\Admin\Documents\HomeLodge-Doc\docs\Presentation\PSM1_Development.pptx"
OUTPUT   = r"c:\Users\Admin\Documents\HomeLodge-Doc\docs\Presentation\HomeLodge_PSM1_v3.pptx"

# Colours from the template
CRIMSON = RGBColor(0x76, 0x1D, 0x3E)
GOLD    = RGBColor(0xFF, 0xC0, 0x00)
BLACK   = RGBColor(0x00, 0x00, 0x00)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
GREY    = RGBColor(0x55, 0x55, 0x55)
DARK    = RGBColor(0x1A, 0x1A, 0x1A)
GREEN   = RGBColor(0x00, 0x70, 0x00)
RED_C   = RGBColor(0xC0, 0x00, 0x00)

prs = Presentation(TEMPLATE)

# ── helpers ────────────────────────────────────────────────────────────────────

def clear_and_set_tf(tf, paragraphs_data, word_wrap=True):
    """
    paragraphs_data: list of tuples:
      (text, font_size, bold, color, level, italic, alignment)
    """
    tf.clear()
    tf.word_wrap = word_wrap
    first = True
    for item in paragraphs_data:
        text      = item[0]
        fs        = item[1] if len(item) > 1 else 16
        bold      = item[2] if len(item) > 2 else False
        color     = item[3] if len(item) > 3 else BLACK
        level     = item[4] if len(item) > 4 else 0
        italic    = item[5] if len(item) > 5 else False
        alignment = item[6] if len(item) > 6 else None
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.level = level
        p.space_before = Pt(3)
        if alignment:
            p.alignment = alignment
        run = p.add_run()
        run.text = text
        run.font.size = Pt(fs)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.italic = italic


def find_shape_by_name(slide, name_part):
    for s in slide.shapes:
        if name_part.lower() in s.name.lower():
            return s
    return None


def get_content_ph(slide):
    for s in slide.shapes:
        if s.shape_type == 14 and "Content" in s.name:
            return s
    return None


def get_title_textbox(slide):
    """Get the big crimson section-title TextBox."""
    for s in slide.shapes:
        if s.shape_type == 17:
            try:
                txt = s.text_frame.text.strip()
                if txt.startswith("(") and ")" in txt:
                    return s
            except:
                pass
    return None


def set_section_title(slide, text, size=40):
    shape = get_title_textbox(slide)
    if shape:
        tf = shape.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = False
        run.font.color.rgb = CRIMSON


def get_gold_textbox(slide):
    """Get the gold-coloured sub-header TextBox."""
    for s in slide.shapes:
        if s.shape_type == 17:
            try:
                for para in s.text_frame.paragraphs:
                    for run in para.runs:
                        try:
                            if run.font.color.rgb == RGBColor(0xFF, 0xC0, 0x00):
                                return s
                        except:
                            pass
            except:
                pass
    return None


def set_gold_subtitle(slide, text, size=28):
    shape = get_gold_textbox(slide)
    if shape:
        tf = shape.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = True
        run.font.color.rgb = GOLD


def fix_textbox_containing(slide, old_text_part, new_text, color=GOLD, size=28, bold=True):
    """Find any TextBox containing old_text_part (case-insensitive) and replace."""
    for s in slide.shapes:
        if s.shape_type == 17 and s.has_text_frame:
            txt = s.text_frame.text.strip()
            if old_text_part.lower() in txt.lower():
                tf = s.text_frame
                tf.clear()
                p = tf.paragraphs[0]
                run = p.add_run()
                run.text = new_text
                run.font.size = Pt(size)
                run.font.bold = bold
                run.font.color.rgb = color
                return True
    return False


def set_content(slide, paragraphs_data):
    shape = get_content_ph(slide)
    if shape:
        clear_and_set_tf(shape.text_frame, paragraphs_data)
        return True
    return False


def add_table_to_slide(slide, rows_data, x, y, w, h, col_widths=None):
    """Add a styled table."""
    n_rows = len(rows_data)
    n_cols = len(rows_data[0])
    tbl_shape = slide.shapes.add_table(n_rows, n_cols, x, y, w, h)
    tbl = tbl_shape.table
    if col_widths:
        for i, cw in enumerate(col_widths):
            tbl.columns[i].width = cw
    for ri, row_data in enumerate(rows_data):
        for ci, cell_text in enumerate(row_data):
            cell = tbl.cell(ri, ci)
            cell.text = cell_text
            tf = cell.text_frame
            tf.word_wrap = True
            para = tf.paragraphs[0]
            para.alignment = PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT
            run = para.runs[0] if para.runs else para.add_run()
            run.text = cell_text
            run.font.size = Pt(12) if ri > 0 else Pt(13)
            run.font.bold = (ri == 0)
            # colours
            if ri == 0:
                run.font.color.rgb = WHITE
            elif ci == n_cols - 1 and cell_text in ("Yes", "✓"):
                run.font.color.rgb = GREEN
                run.font.bold = True
            elif cell_text in ("No", "✗"):
                run.font.color.rgb = RED_C
            else:
                run.font.color.rgb = BLACK
            # cell background
            tc   = cell._tc
            tcPr = tc.get_or_add_tcPr()
            for old in tcPr.findall(qn('a:solidFill')):
                tcPr.remove(old)
            if ri == 0:
                val = "761D3E"
            elif ri % 2 == 0:
                val = "F5E6EC"
            else:
                val = "FFFFFF"
            fill_el = parse_xml(
                f'<a:solidFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
                f'<a:srgbClr val="{val}"/></a:solidFill>'
            )
            tcPr.insert(0, fill_el)
    return tbl_shape


# ═══════════════════════════════════════════════════════════════════════════════
# The template has 31 slides. We select only the ones we need for a 10-slide
# PSM 1 presentation following the format guidelines.
#
# PSM1 Format (10 slides):
#   1. Pendahuluan Pembentangan (Title)
#   2. Pengenalan Projek (Problem Background / Introduction)
#   3. Matlamat Projek (Aim)
#   4. Objektif Projek (Objectives)
#   5. Kajian Latarbelakang Masalah (Lit Review — comparison table)
#   6. Metodologi Pembangunan
#   7. Spesifikasi Keperluan/Input/Output (HW + SW + Tech Stack)
#   8. Hasil Perlaksanaan PSM I dan Perancangan PSM II
#   9. Penutup (Conclusion)
#  10. Terima Kasih / Sesi Soal-Jawab
# ═══════════════════════════════════════════════════════════════════════════════

# Indices of slides to KEEP (0-based), in order
# We pick one slide per section where possible, plus extras for lit review
KEEP_INDICES = [
    0,   # title slide
    3,   # introduction / problem → Slide 2
    5,   # aim → Slide 3
    6,   # objectives → Slide 4
    8,   # lit review A (Airbnb) → Slide 5a
    10,  # lit review B (Booking.com) → Slide 5b
    12,  # lit review C (Agoda) → Slide 5c
    14,  # comparison table → Slide 5d
    18,  # methodology → Slide 6
    19,  # HW requirements → Slide 7a
    15,  # technologies → Slide 7b (SW stack)
    21,  # initial finding (use case) → Slide 8a
    22,  # system design → Slide 8b
    28,  # PSM2 planning → Slide 8c
    29,  # conclusion → Slide 9
    30,  # thank you → Slide 10
]

# Delete slides not in KEEP_INDICES
sldIdLst = prs.slides._sldIdLst
all_ids  = list(sldIdLst)

keep_set = set(KEEP_INDICES)
to_remove = [el for i, el in enumerate(all_ids) if i not in keep_set]
for el in to_remove:
    sldIdLst.remove(el)

print(f"Slides remaining after pruning: {len(prs.slides)}")

# Now reference them in order
slides = list(prs.slides)
S = {name: slides[i] for i, name in enumerate([
    "title", "intro", "aim", "objectives",
    "litA", "litB", "litC", "comparison",
    "method", "hw", "tech",
    "usecase", "design", "psm2",
    "conclusion", "thankyou"
])}


# ── SLIDE 1: TITLE ────────────────────────────────────────────────────────────
slide = S["title"]
for shape in slide.shapes:
    if not shape.has_text_frame:
        continue
    txt = shape.text_frame.text.strip()
    if "HOMESTAY" in txt or "HomeLodge" in txt:
        clear_and_set_tf(shape.text_frame, [
            ("HOMESTAY BOOKING MANAGEMENT SYSTEM", 26, True, CRIMSON, 0),
            ("(HomeLodge)", 26, True, CRIMSON, 0),
        ])
    elif "Aisyah" in txt:
        clear_and_set_tf(shape.text_frame, [
            ("by:", 22, False, BLACK, 0),
            ("Aisyah Umairah binti Azmir", 22, False, CRIMSON, 0),
        ])
    elif "Supervisor" in txt or "Zalmiyah" in txt:
        clear_and_set_tf(shape.text_frame, [
            ("Supervisor:", 22, False, BLACK, 0),
            ("Dr. Zalmiyah binti Zakaria", 22, False, CRIMSON, 0),
        ])
    elif "July 2021" in txt:
        clear_and_set_tf(shape.text_frame, [
            ("22 June 2026", 22, False, BLACK, 0),
        ])
print("✓ Slide 1: Title")


# ── SLIDE 2: PENGENALAN PROJEK ─────────────────────────────────────────────────
slide = S["intro"]
set_section_title(slide, "(1) Pengenalan Projek")
set_content(slide, [
    ("Latarbelakang Masalah:", 18, True, CRIMSON, 0),
    ("Pengendali homestay yang menguruskan pelbagai unit bergantung kepada kaedah manual:", 15, False, GREY, 0),
    ("", 6, False, BLACK, 0),
    ("Tempahan Berganda", 16, True, CRIMSON, 0),
    ("Tiada rekod ketersediaan dikongsi — dua tetamu menempah tarikh yang sama.", 14, False, DARK, 1),
    ("Kutipan Bayaran Manual", 16, True, CRIMSON, 0),
    ("Tiada resit dan tiada proses bayaran balik yang konsisten.", 14, False, DARK, 1),
    ("Penyerahan Kunci Fizikal", 16, True, CRIMSON, 0),
    ("Pengendali perlu hadir secara fizikal setiap kali daftar masuk.", 14, False, DARK, 1),
    ("Komunikasi Tersebar", 16, True, CRIMSON, 0),
    ("Pengesahan dan arahan bayaran terkubur dalam chat peribadi.", 14, False, DARK, 1),
    ("Peringatan Manual", 16, True, CRIMSON, 0),
    ("Tempahan belum dibayar boleh terlepas pandang.", 14, False, DARK, 1),
])
print("✓ Slide 2: Pengenalan Projek")


# ── SLIDE 3: MATLAMAT PROJEK ──────────────────────────────────────────────────
slide = S["aim"]
set_section_title(slide, "(2) Matlamat Projek")
set_content(slide, [
    ("Membangunkan HomeLodge — sistem pengurusan tempahan homestay berasaskan web yang mendigitalkan proses:", 17, False, BLACK, 0),
    ("", 8, False, BLACK, 0),
    ("Tempahan dan ketersediaan", 16, False, DARK, 1),
    ("Pembayaran dalam talian dengan bil dan resit automatik", 16, False, DARK, 1),
    ("Akses pintu melalui QR code (tanpa penyerahan kunci fizikal)", 16, False, DARK, 1),
    ("Chat masa nyata antara tetamu dan admin melalui WebSocket", 16, False, DARK, 1),
    ("Kawalan akses berasaskan peranan (RBAC) dan log audit", 16, False, DARK, 1),
    ("", 8, False, BLACK, 0),
    ("HomeLodge bukan marketplace.", 17, True, CRIMSON, 0),
    ("Alat pengurusan peribadi — tiada komisen, kawalan penuh, penyenaraian peribadi.", 15, False, GREY, 0),
])
print("✓ Slide 3: Matlamat Projek")


# ── SLIDE 4: OBJEKTIF PROJEK ─────────────────────────────────────────────────
slide = S["objectives"]
set_section_title(slide, "(3) Objektif Projek")
set_content(slide, [
    ("01 — Kaji & Analisis", 16, True, CRIMSON, 0),
    ("Mengumpul dan mendokumentasikan semua keperluan fungsian dan bukan fungsian.", 15, False, DARK, 1),
    ("", 6, False, BLACK, 0),
    ("02 — Reka Bentuk", 16, True, CRIMSON, 0),
    ("Menghasilkan seni bina sistem, skema pangkalan data, dan reka bentuk UI.", 15, False, DARK, 1),
    ("", 6, False, BLACK, 0),
    ("03 — Pembangunan", 16, True, CRIMSON, 0),
    ("Membangunkan aplikasi web HomeLodge dengan 11 modul fungsian.", 15, False, DARK, 1),
    ("", 6, False, BLACK, 0),
    ("04 — Pengujian", 16, True, CRIMSON, 0),
    ("Mengesahkan ketepatan, keselamatan, dan kebolehgunaan.", 15, False, DARK, 1),
    ("", 8, False, BLACK, 0),
    ("Skop: Aplikasi web | 2 peranan (Guest, Admin) | Bahasa Inggeris | 11 modul", 14, True, GREY, 0),
    ("Authentication · Homestay · Booking · Payment · Notification · Chat ·", 13, False, GREY, 1),
    ("User Mgmt · Role & Permission · Settings · Audit Log · QR Code Door Access", 13, False, GREY, 1),
])
print("✓ Slide 4: Objektif Projek")


# ── SLIDE 5a: LIT REVIEW – AIRBNB ────────────────────────────────────────────
slide = S["litA"]
set_section_title(slide, "(4) Kajian Latarbelakang")
set_gold_subtitle(slide, "(a) Airbnb", 28)
set_content(slide, [
    ("Marketplace penginapan dalam talian (ditubuhkan 2008) yang menghubungkan hos dengan tetamu.", 16, False, DARK, 0),
    ("", 8, False, BLACK, 0),
    ("Kelebihan:", 15, True, CRIMSON, 0),
    ("Kalendar ketersediaan, pembayaran dalam talian, pemesejan, dashboard hos", 14, False, DARK, 1),
    ("", 6, False, BLACK, 0),
    ("Kekurangan untuk pengendali bebas:", 15, True, CRIMSON, 0),
    ("Tiada akses pintu QR code — penyerahan kunci fizikal masih diperlukan", 14, False, DARK, 1),
    ("Tiada pengurusan lanjutan tempahan dengan aliran pembayaran", 14, False, DARK, 1),
    ("Penyenaraian marketplace wajib — tidak boleh menjalankan tempahan peribadi", 14, False, DARK, 1),
    ("Komisen per-tempahan dikenakan pada setiap transaksi", 14, False, DARK, 1),
    ("Tiada RBAC atau log audit kekal", 14, False, DARK, 1),
])
print("✓ Slide 5a: Lit Review — Airbnb")


# ── SLIDE 5b: LIT REVIEW – BOOKING.COM ───────────────────────────────────────
slide = S["litB"]
set_section_title(slide, "(4) Kajian Latarbelakang")
set_gold_subtitle(slide, "(b) Booking.com", 28)
set_content(slide, [
    ("OTA global digunakan oleh penyedia penginapan dari hotel besar hingga rumah tamu kecil.", 16, False, DARK, 0),
    ("", 8, False, BLACK, 0),
    ("Kelebihan:", 15, True, CRIMSON, 0),
    ("Pengurusan ketersediaan, paparan tempahan, komunikasi tetamu, analitik, program kesetiaan", 14, False, DARK, 1),
    ("", 6, False, BLACK, 0),
    ("Kekurangan untuk pengendali bebas:", 15, True, CRIMSON, 0),
    ("Tiada akses pintu QR code atau kebenaran akses berasaskan peranan", 14, False, DARK, 1),
    ("Tiada pengurusan lanjutan tempahan", 14, False, DARK, 1),
    ("Pengendali terikat dengan struktur komisen dan polisi kandungan platform", 14, False, DARK, 1),
    ("Marketplace memberikan keterlihatan platform ke semua transaksi — bukan alat peribadi", 14, False, DARK, 1),
])
print("✓ Slide 5b: Lit Review — Booking.com")


# ── SLIDE 5c: LIT REVIEW – AGODA ─────────────────────────────────────────────
slide = S["litC"]
set_section_title(slide, "(4) Kajian Latarbelakang")
set_gold_subtitle(slide, "(c) Agoda", 28)
set_content(slide, [
    ("OTA berfokus Asia Tenggara (sebahagian Booking Holdings), popular di Malaysia.", 16, False, DARK, 0),
    ("", 8, False, BLACK, 0),
    ("Kelebihan:", 15, True, CRIMSON, 0),
    ("Penyetempatan pasaran Malaysia: Bahasa Malaysia, pemindahan bank FPX, pemesejan dalam platform", 14, False, DARK, 1),
    ("", 6, False, BLACK, 0),
    ("Kekurangan untuk pengendali bebas:", 15, True, CRIMSON, 0),
    ("Tiada ciri akses pintu QR code", 14, False, DARK, 1),
    ("Tiada sistem pengurusan peranan dan kebenaran", 14, False, DARK, 1),
    ("Tiada aliran kerja lanjutan tempahan formal", 14, False, DARK, 1),
    ("Model marketplace yang sama — berasaskan komisen, penyenaraian awam", 14, False, DARK, 1),
])
# Fallback: fix "(c) System 3" text that isn't gold-coloured in template
fix_textbox_containing(slide, "System 3", "(c) Agoda", GOLD, 28, True)
print("✓ Slide 5c: Lit Review — Agoda")


# ── SLIDE 5d: COMPARISON TABLE ────────────────────────────────────────────────
slide = S["comparison"]
set_section_title(slide, "(4) Kajian Latarbelakang")

# Remove old table and caption
for shape in list(slide.shapes):
    if shape.shape_type == 19:  # TABLE
        shape._element.getparent().remove(shape._element)
    elif shape.has_text_frame and "Table: System comparison" in shape.text_frame.text:
        shape._element.getparent().remove(shape._element)

# Update the caption placeholder if it exists
for shape in slide.shapes:
    if shape.has_text_frame and shape.shape_type == 14 and "Content Placeholder 4" in shape.name:
        clear_and_set_tf(shape.text_frame, [
            ("Jadual: Perbandingan Sistem", 13, False, DARK, 0),
        ])

rows_data = [
    ["Ciri / Feature", "Airbnb", "Booking.com", "Agoda", "HomeLodge"],
    ["Kalendar ketersediaan",           "✓", "✓", "✓", "✓"],
    ["Pembayaran dalam talian",         "✓", "✓", "✓", "✓"],
    ["Pemesejan tetamu-ke-hos",         "✓", "✓", "✓", "✓"],
    ["Akses pintu QR code",             "✗", "✗", "✗", "✓"],
    ["Lanjutan tempahan",               "✗", "✗", "✗", "✓"],
    ["RBAC (kawalan akses peranan)",    "✗", "✗", "✗", "✓"],
    ["Operasi peribadi",                "✗", "✗", "✗", "✓"],
    ["Tiada komisen per-tempahan",      "✗", "✗", "✗", "✓"],
    ["Log audit kekal",                 "✗", "✗", "✗", "✓"],
]
col_widths = [Inches(3.5), Inches(1.4), Inches(1.4), Inches(1.4), Inches(1.5)]
add_table_to_slide(slide, rows_data,
    Inches(0.28), Inches(1.85), Inches(9.44), Inches(5.25), col_widths)
print("✓ Slide 5d: Comparison Table")


# ── SLIDE 6: METODOLOGI PEMBANGUNAN ───────────────────────────────────────────
slide = S["method"]
set_section_title(slide, "(5) Metodologi Pembangunan")
set_content(slide, [
    ("Pendekatan Hibrid: Waterfall + Agile", 17, True, CRIMSON, 0),
    ("", 6, False, BLACK, 0),
    ("Waterfall — Perancangan & Analisis:", 16, True, CRIMSON, 0),
    ("Keperluan fungsian dan bukan fungsian dikumpul dan ditetapkan sebelum pengkodan (URS + PRD)", 14, False, DARK, 1),
    ("Reka bentuk sistem lengkap (ER diagram, use cases, wireframe UI) sebelum pembangunan", 14, False, DARK, 1),
    ("", 6, False, BLACK, 0),
    ("Agile (Scrum + Kanban) — Pembangunan & Pengujian:", 16, True, CRIMSON, 0),
    ("11 modul dibina dalam sprint berperingkat mengikut urutan kebergantungan", 14, False, DARK, 1),
    ("Kanban Board: To Do  →  In Progress  →  In Review  →  Done", 14, False, DARK, 1),
    ("Setiap modul diuji terhadap kriteria penerimaan selepas siap", 14, False, DARK, 1),
    ("", 6, False, BLACK, 0),
    ("Waterfall — Pelancaran:", 16, True, CRIMSON, 0),
    ("Senarai semak berurutan: Nginx, PHP-FPM, MySQL, Redis, SSL (Let's Encrypt)", 14, False, DARK, 1),
    ("", 8, False, BLACK, 0),
    ("Mengapa hibrid? Waterfall untuk dokumentasi formal PSM. Agile untuk 11 modul yang boleh dijejak.", 13, True, GREY, 0),
])
# Clear leftover "Technologies:" gold textbox from reused template slot
fix_textbox_containing(slide, "Technologies", "", GOLD, 28, True)
print("✓ Slide 6: Metodologi")


# ── SLIDE 7a: SPESIFIKASI PERKAKASAN ──────────────────────────────────────────
slide = S["hw"]
set_section_title(slide, "(6) Spesifikasi Keperluan")

# Remove old table
for shape in list(slide.shapes):
    if shape.shape_type == 19:  # TABLE
        shape._element.getparent().remove(shape._element)

# Update caption
for shape in slide.shapes:
    if shape.has_text_frame and shape.shape_type == 14 and "Content Placeholder 4" in shape.name:
        clear_and_set_tf(shape.text_frame, [
            ("Jadual: Spesifikasi Perkakasan", 13, False, DARK, 0),
        ])

hw_rows = [
    ["Komponen", "Spesifikasi Pembangunan", "Spesifikasi Pelayan"],
    ["Pemproses", "Intel Core i5 atau setaraf", "2 vCPU minimum"],
    ["RAM", "8 GB", "4 GB (min 2 GB)"],
    ["Storan", "256 GB SSD", "20 GB SSD minimum"],
    ["Internet", "Broadband, min 10 Mbps", "—"],
    ["Peranti Tambahan", "—", "Smart Lock serasi QR code"],
]
hw_col_widths = [Inches(2.5), Inches(3.2), Inches(3.2)]
add_table_to_slide(slide, hw_rows,
    Inches(0.5), Inches(2.0), Inches(8.9), Inches(4.0), hw_col_widths)
# Clear leftover "Text" content placeholder
for shape in slide.shapes:
    if shape.has_text_frame and shape.shape_type == 14 and shape.text_frame.text.strip() == "Text":
        shape.text_frame.clear()
print("✓ Slide 7a: Spesifikasi Perkakasan")


# ── SLIDE 7b: TEKNOLOGI (PERISIAN) ────────────────────────────────────────────
slide = S["tech"]
set_section_title(slide, "(6) Spesifikasi Keperluan")
set_gold_subtitle(slide, "Stack Teknologi:", 28)

# Remove leftover table from this template slot (it was an SW/HW requirements table)
for shape in list(slide.shapes):
    if shape.shape_type == 19:  # TABLE
        shape._element.getparent().remove(shape._element)
# Clear leftover caption text
for shape in slide.shapes:
    if shape.has_text_frame and shape.shape_type == 14 and ("Software" in shape.text_frame.text or "Table:" in shape.text_frame.text):
        shape.text_frame.clear()
set_content(slide, [
    ("Backend:    Laravel 11 (PHP 8.2) — MVC, ORM, queues, scheduler", 15, False, DARK, 0),
    ("Frontend:   Blade (server-rendered) + Alpine.js + Bootstrap 5", 15, False, DARK, 0),
    ("Database:   MySQL 8 + Eloquent ORM — keutuhan relasional, transaksi atomik", 15, False, DARK, 0),
    ("Real-time:  Laravel Reverb (WebSocket) + Laravel Echo (klien JS)", 15, False, DARK, 0),
    ("Auth:       laravel/breeze + laravel/socialite (Google OAuth 2.0 SSO)", 15, False, DARK, 0),
    ("RBAC:       spatie/laravel-permission — peranan/kebenaran berasaskan pangkalan data", 15, False, DARK, 0),
    ("QR Code:    simplesoftwareio/simple-qrcode — penjanaan lokal, tiada API luaran", 15, False, DARK, 0),
    ("Payment:    Billplz / Toyyibpay / Stripe — webhook-driven, sokongan FPX", 15, False, DARK, 0),
    ("Audit Log:  spatie/laravel-activitylog — rekod peristiwa kekal", 15, False, DARK, 0),
    ("PDF:        barryvdh/laravel-dompdf — bil, resit, eksport laporan", 15, False, DARK, 0),
])
print("✓ Slide 7b: Stack Teknologi")


# ── SLIDE 8a: USE CASE MODEL ──────────────────────────────────────────────────
slide = S["usecase"]
set_section_title(slide, "(7) Hasil Perlaksanaan PSM I")
set_content(slide, [
    ("Model Use-Case:", 17, True, CRIMSON, 0),
    ("", 6, False, BLACK, 0),
    ("4 Aktor:", 15, True, CRIMSON, 0),
    ("Guest (Manusia) — layari, tempah, bayar, terima QR code, chat, beri maklum balas", 14, False, DARK, 1),
    ("Admin (Manusia) — urus unit, tempahan, pengguna, peranan, tetapan, laporan", 14, False, DARK, 1),
    ("Payment Gateway (Luaran) — proses pembayaran, hantar webhook callbacks", 14, False, DARK, 1),
    ("Google Calendar API (Luaran) — terima peristiwa tempahan selepas pengesahan", 14, False, DARK, 1),
    ("", 6, False, BLACK, 0),
    ("12 Modul  |  47 Use Cases  |  Semua dijejak kepada keperluan URS dan PRD", 16, True, CRIMSON, 0),
    ("Authentication, Homestay, Booking, Payment, Notification, Chat", 14, False, DARK, 1),
    ("User Mgmt, Role & Permission, Settings, Audit Logs, QR Code, Reporting", 14, False, DARK, 1),
    ("", 6, False, BLACK, 0),
    ("Aliran utama: Tempahan tetamu, Pembayaran, Auto-pembatalan, Kitaran QR Code, Lanjutan tempahan", 13, False, GREY, 0),
])
print("✓ Slide 8a: Use Case Model")


# ── SLIDE 8b: SYSTEM DESIGN ──────────────────────────────────────────────────
slide = S["design"]
set_section_title(slide, "(7) Hasil Perlaksanaan PSM I")
set_gold_subtitle(slide, "Reka Bentuk Sistem:", 28)
set_content(slide, [
    ("Seni Bina (MVC + Service Layer):", 15, True, CRIMSON, 0),
    ("Client: Blade + Alpine.js + Bootstrap 5 + Laravel Echo", 14, False, DARK, 1),
    ("Server: Nginx → PHP-FPM → Laravel 11 | Queue Worker | Scheduler | Reverb", 14, False, DARK, 1),
    ("Data: MySQL 8 (primer) | Redis 7 (cache, queues, sessions)", 14, False, DARK, 1),
    ("External: Payment Gateway | Google OAuth + Calendar | SMTP", 14, False, DARK, 1),
    ("", 6, False, BLACK, 0),
    ("Reka Bentuk Pangkalan Data:", 15, True, CRIMSON, 0),
    ("17 entiti | MySQL 8 + Eloquent ORM | FK constraints + soft deletes", 14, False, DARK, 1),
    ("Jadual utama: bookings, booking_extensions, payments, qr_codes, bills", 14, False, DARK, 1),
    ("settings (key-value model — tiada migrasi skema untuk konfigurasi baru)", 14, False, DARK, 1),
    ("Webhook idempotent melalui indeks gateway_reference unik", 14, False, DARK, 1),
    ("", 6, False, BLACK, 0),
    ("Reka Bentuk Antara Muka:", 15, True, CRIMSON, 0),
    ("Guest: top navbar, aliran linear (layari → tempah → bayar → QR code)", 14, False, DARK, 1),
    ("Admin: sidebar kiri tetap, 10 kumpulan navigasi, responsif (mobile/tablet/desktop)", 14, False, DARK, 1),
])
print("✓ Slide 8b: System Design")


# ── SLIDE 8c: PSM 2 PLANNING ─────────────────────────────────────────────────
slide = S["psm2"]
set_section_title(slide, "(8) Perancangan PSM II")
set_content(slide, [
    ("PSM II memfokuskan sepenuhnya kepada pembangunan, pengujian, dan pelancaran sistem.", 16, False, GREY, 0),
    ("", 6, False, BLACK, 0),
    ("Fasa 1 — Persediaan Persekitaran", 15, True, CRIMSON, 0),
    ("Konfigurasi Laragon, MySQL, Redis, dan semua dependensi Composer + npm", 14, False, DARK, 1),
    ("", 4, False, BLACK, 0),
    ("Fasa 2 — Pembangunan Backend (Sprint Agile)", 15, True, CRIMSON, 0),
    ("Bina semua 11 modul mengikut urutan kebergantungan menggunakan Laravel 11", 14, False, DARK, 1),
    ("Integrasi payment gateway (webhooks), Laravel Reverb (WebSocket), penjanaan QR code", 14, False, DARK, 1),
    ("", 4, False, BLACK, 0),
    ("Fasa 3 — Pembangunan Frontend (serentak dengan backend)", 15, True, CRIMSON, 0),
    ("Pelaksanaan wireframe sebagai halaman Blade + Alpine.js dengan Bootstrap 5", 14, False, DARK, 1),
    ("", 4, False, BLACK, 0),
    ("Fasa 4 — Pengujian", 15, True, CRIMSON, 0),
    ("Unit test, integration test, dan User Acceptance Testing (UAT)", 14, False, DARK, 1),
    ("", 4, False, BLACK, 0),
    ("Fasa 5 — Pelancaran", 15, True, CRIMSON, 0),
    ("Deploy pada Nginx + PHP-FPM + MySQL 8 + Redis; aktifkan SSL dengan Let's Encrypt", 14, False, DARK, 1),
])
print("✓ Slide 8c: PSM 2 Planning")


# ── SLIDE 9: PENUTUP ──────────────────────────────────────────────────────────
slide = S["conclusion"]
set_section_title(slide, "(9) Penutup")
set_content(slide, [
    ("Rumusan Masalah:", 16, True, CRIMSON, 0),
    ("Tempahan berganda, kutipan bayaran manual, penyerahan kunci fizikal, komunikasi tersebar", 14, False, DARK, 1),
    ("", 6, False, BLACK, 0),
    ("Platform sedia ada tidak mencukupi:", 16, True, CRIMSON, 0),
    ("Airbnb, Booking.com, dan Agoda mengendalikan ketersediaan dan pembayaran tetapi beroperasi sebagai marketplace awam.", 14, False, DARK, 1),
    ("Tiada akses pintu QR, lanjutan tempahan, RBAC halus, atau operasi peribadi tanpa komisen.", 14, False, DARK, 1),
    ("", 6, False, BLACK, 0),
    ("Hasil PSM I yang telah disiapkan:", 16, True, CRIMSON, 0),
    ("Spesifikasi keperluan lengkap (URS + PRD)", 14, False, DARK, 1),
    ("Reka bentuk sistem penuh: seni bina, ER diagram (17 entiti), wireframe UI", 14, False, DARK, 1),
    ("47 use cases merentasi 12 modul didokumentasikan dan dijejak", 14, False, DARK, 1),
    ("Metodologi hibrid Waterfall + Agile dipilih dan dijustifikasikan", 14, False, DARK, 1),
    ("Stack teknologi dinilai dan dipilih", 14, False, DARK, 1),
    ("", 8, False, BLACK, 0),
    ("PSM II akan melaksanakan, menguji, dan melancarkan HomeLodge.", 16, True, CRIMSON, 0),
])
print("✓ Slide 9: Penutup")


# ── SLIDE 10: TERIMA KASIH ───────────────────────────────────────────────────
# Keep template content as-is; it already has "Thank You" and "Q&A"
print("✓ Slide 10: Terima Kasih (template kept)")


# ── SAVE ──────────────────────────────────────────────────────────────────────
prs.save(OUTPUT)
print(f"\n✅ Saved: {OUTPUT}")
print(f"Total slides: {len(prs.slides)}")
