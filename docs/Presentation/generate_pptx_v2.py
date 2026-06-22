"""
HomeLodge PSM I - Presentation v2
Uses PSM1_Development.pptx as the template.
Correct approach: modify the template in-place by replacing slide content.
"""

import copy
import lxml.etree as etree
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml

TEMPLATE = r"c:\Users\Admin\Documents\HomeLodge-Doc\docs\Presentation\PSM1_Development.pptx"
OUTPUT   = r"c:\Users\Admin\Documents\HomeLodge-Doc\docs\Presentation\HomeLodge_PSM1_v2.pptx"

# Colours from the template
CRIMSON = RGBColor(0x76, 0x1D, 0x3E)
GOLD    = RGBColor(0xFF, 0xC0, 0x00)
BLACK   = RGBColor(0x00, 0x00, 0x00)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
GREY    = RGBColor(0x55, 0x55, 0x55)
DARK    = RGBColor(0x1A, 0x1A, 0x1A)
GREEN   = RGBColor(0x00, 0x70, 0x00)
RED_C   = RGBColor(0xC0, 0x00, 0x00)

prs = Presentation(TEMPLATE)

# ── helpers ────────────────────────────────────────────────────────────────────

def clear_and_set_tf(tf, paragraphs_data, word_wrap=True):
    """
    paragraphs_data: list of (text, font_size, bold, color, level, italic)
    """
    tf.clear()
    tf.word_wrap = word_wrap
    first = True
    for item in paragraphs_data:
        text      = item[0]
        fs        = item[1] if len(item) > 1 else 16
        bold      = item[2] if len(item) > 2 else False
        color     = item[3] if len(item) > 3 else BLACK
        level     = item[4] if len(item) > 4 else 0
        italic    = item[5] if len(item) > 5 else False
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.level = level
        p.space_before = Pt(3)
        run = p.add_run()
        run.text = text
        run.font.size = Pt(fs)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.italic = italic

def find_shape_by_name(slide, name_part):
    for s in slide.shapes:
        if name_part.lower() in s.name.lower():
            return s
    return None

def get_content_ph(slide):
    for s in slide.shapes:
        if s.shape_type == 14 and "Content" in s.name:
            return s
    return None

def get_title_textbox(slide):
    """Get the big crimson section-title TextBox."""
    for s in slide.shapes:
        if s.shape_type == 17:
            try:
                txt = s.text_frame.text.strip()
                if txt.startswith("(") and ")" in txt:
                    return s
            except:
                pass
    return None

def set_section_title(slide, text, size=40):
    shape = get_title_textbox(slide)
    if shape:
        tf = shape.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = False
        run.font.color.rgb = CRIMSON

def get_gold_textbox(slide):
    """Get the gold-coloured sub-header TextBox."""
    for s in slide.shapes:
        if s.shape_type == 17:
            try:
                for para in s.text_frame.paragraphs:
                    for run in para.runs:
                        try:
                            if run.font.color.rgb == RGBColor(0xFF, 0xC0, 0x00):
                                return s
                        except:
                            pass
            except:
                pass
    return None

def set_gold_subtitle(slide, text, size=28):
    shape = get_gold_textbox(slide)
    if shape:
        tf = shape.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = True
        run.font.color.rgb = GOLD

def set_content(slide, paragraphs_data):
    shape = get_content_ph(slide)
    if shape:
        clear_and_set_tf(shape.text_frame, paragraphs_data)
        return True
    return False

def add_table_to_slide(slide, rows_data, x, y, w, h, col_widths=None):
    """Add a styled table."""
    n_rows = len(rows_data)
    n_cols = len(rows_data[0])
    tbl_shape = slide.shapes.add_table(n_rows, n_cols, x, y, w, h)
    tbl = tbl_shape.table
    if col_widths:
        for i, cw in enumerate(col_widths):
            tbl.columns[i].width = cw
    for ri, row_data in enumerate(rows_data):
        for ci, cell_text in enumerate(row_data):
            cell = tbl.cell(ri, ci)
            cell.text = cell_text
            tf = cell.text_frame
            tf.word_wrap = True
            para = tf.paragraphs[0]
            para.alignment = PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT
            run = para.runs[0] if para.runs else para.add_run()
            run.text = cell_text
            run.font.size = Pt(12) if ri > 0 else Pt(13)
            run.font.bold = (ri == 0)
            # colours
            if ri == 0:
                run.font.color.rgb = WHITE
            elif ci == n_cols - 1 and cell_text == "Yes":
                run.font.color.rgb = GREEN
                run.font.bold = True
            elif cell_text == "No":
                run.font.color.rgb = RED_C
            else:
                run.font.color.rgb = BLACK
            # cell background
            tc   = cell._tc
            tcPr = tc.get_or_add_tcPr()
            for old in tcPr.findall(qn('a:solidFill')):
                tcPr.remove(old)
            if ri == 0:
                val = "761D3E"
            elif ri % 2 == 0:
                val = "F5E6EC"
            else:
                val = "FFFFFF"
            fill_el = parse_xml(
                f'<a:solidFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
                f'<a:srgbClr val="{val}"/></a:solidFill>'
            )
            tcPr.insert(0, fill_el)
    return tbl_shape

# ═══════════════════════════════════════════════════════════════════════════════
# We will keep only the slides we need.
# The template has 31 slides. We pick:
#   slide[0]  = title slide
#   slide[1]  = outlines
#   slide[3]  = problem background / intro content
#   slide[4]  = second content slide (aim)
#   slide[5]  = third content slide (objectives)
#   slide[7]  = scope
#   slide[8]  = lit review A
#   slide[10] = lit review B
#   slide[12] = lit review C
#   slide[14] = comparison table
#   slide[15] = technologies
#   slide[16] = methodology
#   slide[19] = hardware
#   slide[20] = software
#   slide[22] = initial finding use case
#   slide[23] = system design text
#   slide[27] = system interfaces text
#   slide[28] = PSM2 planning
#   slide[29] = conclusion
#   slide[30] = thank you
#
# Strategy: DELETE all other slides from the presentation, keeping only those
# we need. Then edit each one.
# ═══════════════════════════════════════════════════════════════════════════════

# Indices of slides to KEEP (0-based), in order
KEEP_INDICES = [0, 1, 3, 4, 5, 7, 8, 10, 12, 14, 15, 16, 19, 20, 22, 23, 27, 28, 29, 30]

# Delete slides not in KEEP_INDICES
# We'll do it by removing from the sldIdLst
sldIdLst = prs.slides._sldIdLst
all_ids  = list(sldIdLst)   # list of <p:sldId> elements

keep_set = set(KEEP_INDICES)
to_remove = [el for i, el in enumerate(all_ids) if i not in keep_set]
for el in to_remove:
    sldIdLst.remove(el)

print(f"Slides remaining after pruning: {len(prs.slides)}")

# Now reference them in order
slides = list(prs.slides)
# Map: our logical name -> slides list index
# After pruning, slides are in the original order of KEEP_INDICES
S = {name: slides[i] for i, name in enumerate([
    "title", "outline", "intro", "problem", "aim",
    "scope", "litA", "litB", "litC", "comparison",
    "tech", "method", "hw", "sw", "usecase",
    "dbdesign", "arch", "psm2", "conclusion", "thankyou"
])}

# ── SLIDE: TITLE ───────────────────────────────────────────────────────────────
slide = S["title"]
for shape in slide.shapes:
    if not shape.has_text_frame:
        continue
    txt = shape.text_frame.text.strip()
    if "HOMESTAY" in txt or "HomeLodge" in txt:
        clear_and_set_tf(shape.text_frame, [
            ("HOMESTAY BOOKING MANAGEMENT SYSTEM", 26, True, CRIMSON, 0),
            ("(HomeLodge)", 26, True, CRIMSON, 0),
        ])
    elif txt.startswith("by") or (txt == "" and shape.top > Inches(2.5) and shape.top < Inches(3.5)):
        pass  # keep
    elif "Aisyah" in txt:
        clear_and_set_tf(shape.text_frame, [
            ("by:", 22, False, BLACK, 0),
            ("Aisyah Umairah binti Azmir", 22, False, CRIMSON, 0),
        ])
    elif "Supervisor" in txt or "Zalmiyah" in txt:
        clear_and_set_tf(shape.text_frame, [
            ("Supervisor:", 22, False, BLACK, 0),
            ("Dr. Zalmiyah binti Zakaria", 22, False, CRIMSON, 0),
        ])
    elif "July 2021" in txt:
        clear_and_set_tf(shape.text_frame, [
            ("21 June 2026", 22, False, BLACK, 0),
        ])
print("Title slide done")

# ── SLIDE: OUTLINE ─────────────────────────────────────────────────────────────
slide = S["outline"]
ph = get_content_ph(slide)
if ph:
    clear_and_set_tf(ph.text_frame, [
        ("Introduction",                              20, False, CRIMSON, 0),
        ("Problem Background",                        20, False, CRIMSON, 0),
        ("Project Aim & Objectives",                  20, False, CRIMSON, 0),
        ("Literature Review",                         20, False, CRIMSON, 0),
        ("Methodology",                               20, False, CRIMSON, 0),
        ("System Requirements",                       20, False, CRIMSON, 0),
        ("Initial Finding (Analysis & Design)",       20, False, CRIMSON, 0),
        ("PSM 2 Planning",                            20, False, CRIMSON, 0),
        ("Conclusion",                                20, False, CRIMSON, 0),
    ])
print("Outline slide done")

# ── SLIDE: INTRODUCTION ────────────────────────────────────────────────────────
slide = S["intro"]
set_section_title(slide, "(1) Introduction")
set_content(slide, [
    ("HomeLodge is a web-based homestay booking management system.", 18, False, BLACK, 0),
    ("", 10, False, BLACK, 0),
    ("Designed for operators managing more than one property who currently rely on messaging apps and phone calls with no formal system.", 16, False, GREY, 0),
    ("", 10, False, BLACK, 0),
    ("The system covers:", 16, False, DARK, 0),
    ("Availability checking and booking submission", 15, False, DARK, 1),
    ("Online payment with automated billing and receipts", 15, False, DARK, 1),
    ("Time-limited QR code door access (no physical key handover needed)", 15, False, DARK, 1),
    ("Real-time chat between guest and admin via WebSocket", 15, False, DARK, 1),
    ("Role-based access control and full immutable audit logging", 15, False, DARK, 1),
])
print("Introduction slide done")

# ── SLIDE: PROBLEM BACKGROUND ──────────────────────────────────────────────────
slide = S["problem"]
set_section_title(slide, "(2) Problem Background")
set_content(slide, [
    ("Operators managing multiple units through informal channels face recurring failures:", 16, False, GREY, 0),
    ("", 8, False, BLACK, 0),
    ("Double-bookings", 16, True, CRIMSON, 0),
    ("No shared availability record. Two guests can book the same dates on the same property.", 15, False, DARK, 1),
    ("Unstructured payment collection", 16, True, CRIMSON, 0),
    ("No receipts, no consistent refund process; operators calculate everything manually.", 15, False, DARK, 1),
    ("Physical key handover", 16, True, CRIMSON, 0),
    ("Operator must be physically present at each check-in. Impossible when multiple units have simultaneous arrivals.", 15, False, DARK, 1),
    ("Scattered communication", 16, True, CRIMSON, 0),
    ("Confirmations and payment instructions buried in personal chat threads with no link to booking records.", 15, False, DARK, 1),
    ("Manual reminders", 16, True, CRIMSON, 0),
    ("Easy to miss when managing several units; unpaid bookings can lapse unnoticed.", 15, False, DARK, 1),
])
print("Problem Background slide done")

# ── SLIDE: AIM ─────────────────────────────────────────────────────────────────
slide = S["aim"]
set_section_title(slide, "(3) Project Aim")
set_gold_subtitle(slide, "Project Aim:", 28)
set_content(slide, [
    ("To develop HomeLodge, a web-based homestay booking management system that digitalises the reservation, payment, property access, and communication processes for operators managing multiple homestay units.", 17, False, BLACK, 0),
    ("", 10, False, BLACK, 0),
    ("HomeLodge is not a marketplace.", 17, True, CRIMSON, 0),
    ("It is a private management tool that the operator owns, configures, and controls — with no mandatory public listing and no per-booking commission.", 16, False, DARK, 0),
])
print("Aim slide done")

# ── SLIDE: OBJECTIVES ─────────────────────────────────────────────────────────
slide = S["scope"]  # reusing the 'scope' template slot for objectives, then scope
# Actually we kept the scope slide (index 7) as S["scope"]. Let's use aim slot for objectives
# and scope slot for scope. But we only have one slot between them.
# Let's put objectives content in S["aim"] after setting aim, and scope in S["scope"].
# The "aim" slide is already done above. Now do scope:

slide = S["scope"]
set_section_title(slide, "(3) Project Objectives & Scope")
set_content(slide, [
    ("Objectives:", 16, True, CRIMSON, 0),
    ("1. Study & Analyse — gather and document all functional and non-functional requirements", 15, False, DARK, 1),
    ("2. Design — produce system architecture, database schema, and UI for all modules and both roles", 15, False, DARK, 1),
    ("3. Develop — build HomeLodge with 11 functional modules", 15, False, DARK, 1),
    ("4. Test — verify correctness, security, and usability against stated requirements", 15, False, DARK, 1),
    ("", 8, False, BLACK, 0),
    ("Scope:", 16, True, CRIMSON, 0),
    ("Web-based application; two user roles: Guest and Admin", 15, False, DARK, 1),
    ("Operates in English only; multiple units under one administrator account", 15, False, DARK, 1),
    ("11 modules: Authentication, Homestay Management, Booking, Payment, Notification,", 15, False, DARK, 1),
    ("Chat, User Management, Role & Permission, System Settings, Audit Logs, QR Code Door Access", 15, False, DARK, 1),
    ("QR code module designed for compatible smart lock devices; physical installation is outside scope", 15, False, DARK, 1),
])
print("Objectives & Scope slide done")

# ── SLIDE: LIT REVIEW – AIRBNB ─────────────────────────────────────────────────
slide = S["litA"]
set_section_title(slide, "(4) Literature Review")
set_gold_subtitle(slide, "(a) Airbnb", 28)
set_content(slide, [
    ("Airbnb is an online accommodation marketplace (founded 2008) connecting hosts with guests globally.", 16, False, DARK, 0),
    ("", 8, False, BLACK, 0),
    ("What it provides:", 15, True, CRIMSON, 0),
    ("Centralised availability calendar, online payment processing, in-app messaging, host dashboard", 15, False, DARK, 1),
    ("", 8, False, BLACK, 0),
    ("What it lacks for independent operators:", 15, True, CRIMSON, 0),
    ("No QR code door access — physical key handover still required", 15, False, DARK, 1),
    ("No booking extension management with a payment workflow", 15, False, DARK, 1),
    ("Mandatory public listing — operator cannot run private, direct bookings", 15, False, DARK, 1),
    ("Per-booking commission charged on every transaction", 15, False, DARK, 1),
    ("No role-based access control or immutable audit log", 15, False, DARK, 1),
])
print("Lit Review Airbnb done")

# ── SLIDE: LIT REVIEW – BOOKING.COM ────────────────────────────────────────────
slide = S["litB"]
set_section_title(slide, "(4) Literature Review")
set_gold_subtitle(slide, "(b) Booking.com", 28)
set_content(slide, [
    ("Booking.com is a global OTA used by accommodation providers from large hotels to small guesthouses.", 16, False, DARK, 0),
    ("", 8, False, BLACK, 0),
    ("What it provides:", 15, True, CRIMSON, 0),
    ("Availability management, incoming booking view, guest communication, analytics, loyalty programme", 15, False, DARK, 1),
    ("", 8, False, BLACK, 0),
    ("What it lacks for independent operators:", 15, True, CRIMSON, 0),
    ("No QR code door access or role-based access permissions", 15, False, DARK, 1),
    ("No booking extension management", 15, False, DARK, 1),
    ("Operators bound by platform commission structure and content policies", 15, False, DARK, 1),
    ("Marketplace gives platform visibility into all transactions — not a private tool", 15, False, DARK, 1),
])
print("Lit Review Booking.com done")

# ── SLIDE: LIT REVIEW – AGODA ─────────────────────────────────────────────────
slide = S["litC"]
set_section_title(slide, "(4) Literature Review")
set_gold_subtitle(slide, "(c) Agoda", 28)
set_content(slide, [
    ("Agoda is a Southeast Asia-focused OTA (part of Booking Holdings) widely used in Malaysia.", 16, False, DARK, 0),
    ("", 8, False, BLACK, 0),
    ("What it provides:", 15, True, CRIMSON, 0),
    ("Localisation for Malaysian market: Bahasa Malaysia, FPX bank transfers, in-platform messaging", 15, False, DARK, 1),
    ("", 8, False, BLACK, 0),
    ("What it lacks for independent operators:", 15, True, CRIMSON, 0),
    ("No QR code door access feature", 15, False, DARK, 1),
    ("No role and permission management system", 15, False, DARK, 1),
    ("No formal booking extension workflow", 15, False, DARK, 1),
    ("Same marketplace model as Airbnb and Booking.com — commission-based, public listing", 15, False, DARK, 1),
])
print("Lit Review Agoda done")

# ── SLIDE: COMPARISON TABLE ────────────────────────────────────────────────────
slide = S["comparison"]
set_section_title(slide, "(4) Literature Review")

# Remove the old table and caption shapes
for shape in list(slide.shapes):
    if shape.shape_type == 19:   # TABLE
        shape._element.getparent().remove(shape._element)
    elif shape.has_text_frame and "Table: System comparison" in shape.text_frame.text:
        shape._element.getparent().remove(shape._element)

rows_data = [
    ["Feature", "Airbnb", "Booking.com", "Agoda", "HomeLodge"],
    ["Online availability calendar",    "Yes", "Yes", "Yes", "Yes"],
    ["Online payment processing",       "Yes", "Yes", "Yes", "Yes"],
    ["In-app guest-to-host messaging",  "Yes", "Yes", "Yes", "Yes"],
    ["QR code door access",             "No",  "No",  "No",  "Yes"],
    ["Booking extension management",    "No",  "No",  "No",  "Yes"],
    ["Role-based access control",       "No",  "No",  "No",  "Yes"],
    ["Private, non-marketplace",        "No",  "No",  "No",  "Yes"],
    ["No per-booking commission",       "No",  "No",  "No",  "Yes"],
    ["Immutable audit log",             "No",  "No",  "No",  "Yes"],
]
col_widths = [Inches(3.5), Inches(1.4), Inches(1.4), Inches(1.4), Inches(1.5)]
add_table_to_slide(slide, rows_data,
    Inches(0.28), Inches(1.85), Inches(9.44), Inches(5.25), col_widths)
print("Comparison table done")

# ── SLIDE: TECHNOLOGY STACK ────────────────────────────────────────────────────
slide = S["tech"]
set_section_title(slide, "(4) Literature Review")
set_gold_subtitle(slide, "Technologies:", 28)
set_content(slide, [
    ("Backend:    Laravel 11 (PHP 8.2) — MVC, built-in ORM, queues, scheduler, first-party packages", 15, False, DARK, 0),
    ("Frontend:   Blade (server-rendered) + Alpine.js (client-side state) + Bootstrap 5", 15, False, DARK, 0),
    ("Database:   MySQL 8 + Eloquent ORM — relational integrity, atomic transactions", 15, False, DARK, 0),
    ("Real-time:  Laravel Reverb (self-hosted WebSocket) + Laravel Echo (JS client)", 15, False, DARK, 0),
    ("Auth:       laravel/breeze + laravel/socialite (Google OAuth 2.0 SSO)", 15, False, DARK, 0),
    ("RBAC:       spatie/laravel-permission — database-backed, fine-grained roles/permissions", 15, False, DARK, 0),
    ("QR Code:    simplesoftwareio/simple-qrcode — local generation, no external API dependency", 15, False, DARK, 0),
    ("Payment:    Billplz / Toyyibpay / Stripe — webhook-driven, FPX support for Malaysian market", 15, False, DARK, 0),
    ("Audit Log:  spatie/laravel-activitylog — read-only immutable event records", 15, False, DARK, 0),
    ("PDF:        barryvdh/laravel-dompdf — bills, receipts, report exports", 15, False, DARK, 0),
])
print("Technologies slide done")

# ── SLIDE: METHODOLOGY ─────────────────────────────────────────────────────────
slide = S["method"]
set_section_title(slide, "(5) Methodology")
set_gold_subtitle(slide, "Hybrid: Waterfall + Agile", 28)
set_content(slide, [
    ("Waterfall — Planning & Analysis:", 16, True, CRIMSON, 0),
    ("All functional and non-functional requirements gathered and frozen before coding (URS + PRD)", 15, False, DARK, 1),
    ("Complete system design (ER diagram, use cases, UI wireframes) before implementation", 15, False, DARK, 1),
    ("", 8, False, BLACK, 0),
    ("Agile (Scrum + Kanban Board) — Implementation & Testing:", 16, True, CRIMSON, 0),
    ("11 modules built in time-boxed sprints in dependency order (e.g. Payment after Booking)", 15, False, DARK, 1),
    ("Kanban Board states: To Do  ->  In Progress  ->  In Review  ->  Done", 15, False, DARK, 1),
    ("Each module tested against acceptance criteria immediately after completion", 15, False, DARK, 1),
    ("", 8, False, BLACK, 0),
    ("Waterfall — Deployment:", 16, True, CRIMSON, 0),
    ("Sequential checklist: Nginx, PHP-FPM, MySQL, Redis, SSL (Let's Encrypt)", 15, False, DARK, 1),
    ("", 8, False, BLACK, 0),
    ("Why hybrid? Waterfall gives the formal PSM documentation. Agile makes 11 modules trackable.", 14, True, GREY, 0),
])
print("Methodology slide done")

# ── SLIDE: HARDWARE REQUIREMENTS ──────────────────────────────────────────────
slide = S["hw"]
set_section_title(slide, "(6) System Requirements")
set_gold_subtitle(slide, "Hardware Specification:", 28)
set_content(slide, [
    ("Development Hardware:", 16, True, CRIMSON, 0),
    ("Laptop / Desktop — Intel Core i5 or equivalent, 8 GB RAM, 256 GB SSD", 15, False, DARK, 1),
    ("Internet — Broadband, minimum 10 Mbps", 15, False, DARK, 1),
    ("", 8, False, BLACK, 0),
    ("Deployment Server:", 16, True, CRIMSON, 0),
    ("CPU — 2 vCPU minimum (Nginx, PHP-FPM, Queue Workers, Reverb WebSocket process)", 15, False, DARK, 1),
    ("RAM — 2 GB minimum, 4 GB recommended (Redis cache, queue, session)", 15, False, DARK, 1),
    ("Storage — 20 GB SSD minimum (database, media uploads, logs)", 15, False, DARK, 1),
    ("Smart Lock Device — QR-code-compatible; physical lock installation is outside project scope", 15, False, DARK, 1),
])
print("Hardware slide done")

# ── SLIDE: SOFTWARE REQUIREMENTS ──────────────────────────────────────────────
slide = S["sw"]
set_section_title(slide, "(6) System Requirements")
set_gold_subtitle(slide, "Software Specification:", 28)
set_content(slide, [
    ("Development Software:", 16, True, CRIMSON, 0),
    ("PHP 8.2  |  Composer  |  Node.js 18+ / npm  |  Git", 15, False, DARK, 1),
    ("Laragon (local dev environment — bundles Nginx, MySQL, PHP for Windows)", 15, False, DARK, 1),
    ("MySQL 8.x  |  Redis 7.x  |  Visual Studio Code / PhpStorm  |  Postman  |  TablePlus", 15, False, DARK, 1),
    ("", 8, False, BLACK, 0),
    ("Deployment Software:", 16, True, CRIMSON, 0),
    ("Ubuntu Server 22.04 LTS  |  Nginx  |  PHP-FPM 8.2  |  Composer", 15, False, DARK, 1),
    ("MySQL 8.x  |  Redis 7.x  |  Node.js 18+ (for npm run build)", 15, False, DARK, 1),
    ("Let's Encrypt / Certbot (SSL/TLS certificate)", 15, False, DARK, 1),
])
print("Software slide done")

# ── SLIDE: USE CASE MODEL ──────────────────────────────────────────────────────
slide = S["usecase"]
set_section_title(slide, "(7) Initial Finding")
set_gold_subtitle(slide, "Use Case Model:", 28)
set_content(slide, [
    ("4 Actors:", 16, True, CRIMSON, 0),
    ("Guest (Human) — browses, books, pays, receives QR code, chats, leaves feedback", 15, False, DARK, 1),
    ("Admin (Human) — manages units, bookings, users, roles, settings, reports", 15, False, DARK, 1),
    ("Payment Gateway (External) — processes payments, sends webhook callbacks to HomeLodge", 15, False, DARK, 1),
    ("Google Calendar API (External) — receives booking events upon confirmation", 15, False, DARK, 1),
    ("", 8, False, BLACK, 0),
    ("12 Modules  |  47 Use Cases  |  All traced to URS and PRD requirements", 16, True, CRIMSON, 0),
    ("Authentication, Homestay Management, Booking, Payment, Notification, Chat", 15, False, DARK, 1),
    ("User Management, Role & Permission, System Settings, Audit Logs, QR Code, Reporting & Analytics", 15, False, DARK, 1),
    ("", 8, False, BLACK, 0),
    ("Key workflows: Guest booking flow, Payment flow, Auto-cancellation, QR Code lifecycle, Booking extension flow", 14, False, GREY, 0),
])
print("Use Case slide done")

# ── SLIDE: DATABASE DESIGN ────────────────────────────────────────────────────
slide = S["dbdesign"]
set_section_title(slide, "(7) Initial Finding")
set_gold_subtitle(slide, "Database Design:", 28)
set_content(slide, [
    ("16 tables  |  MySQL 8 + Eloquent ORM  |  Foreign key constraints + soft deletes", 16, True, CRIMSON, 0),
    ("", 8, False, BLACK, 0),
    ("Core tables:", 15, True, DARK, 0),
    ("users — auth fields, lockout tracking, Google SSO id, must_change_password flag", 15, False, DARK, 1),
    ("homestays — name, pricing, check-in/out times, per-unit extension payment window", 15, False, DARK, 1),
    ("bookings — central reservation; status lifecycle: pending_payment -> confirmed -> completed/cancelled", 15, False, DARK, 1),
    ("booking_extensions — stores original + extended dates, extra charge, payment deadline", 15, False, DARK, 1),
    ("payments / bills / refunds — full billing cycle with gateway reference for idempotent webhooks", 15, False, DARK, 1),
    ("qr_codes — unique token per booking, valid_from / valid_until timestamps, status", 15, False, DARK, 1),
    ("settings — key-value config (SMTP, security, extension rates, refund policy)", 15, False, DARK, 1),
    ("", 8, False, BLACK, 0),
    ("ER Diagram and Data Dictionary are included in the full report (Chapter 4).", 13, False, GREY, 0),
])
print("DB Design slide done")

# ── SLIDE: SYSTEM ARCHITECTURE ────────────────────────────────────────────────
slide = S["arch"]
set_section_title(slide, "(7) Initial Finding")
set_gold_subtitle(slide, "System Architecture (MVC):", 28)
set_content(slide, [
    ("Client (Browser):", 15, True, CRIMSON, 0),
    ("Blade (server-rendered HTML)  |  Alpine.js (client-side state)  |  Bootstrap 5  |  Laravel Echo (WebSocket)", 14, False, DARK, 1),
    ("", 6, False, BLACK, 0),
    ("Application Server:", 15, True, CRIMSON, 0),
    ("Nginx  ->  PHP-FPM  ->  Laravel 11  |  Queue Worker  |  Laravel Scheduler (cron)  |  Reverb (WebSocket server)", 14, False, DARK, 1),
    ("", 6, False, BLACK, 0),
    ("Scheduled background jobs:", 15, True, CRIMSON, 0),
    ("Auto-cancel unpaid bookings (hourly)  |  QR code expiry (at checkout time)", 14, False, DARK, 1),
    ("Payment reminders (daily)  |  Stay reminders (daily)  |  Extension revert (every 1-5 min)", 14, False, DARK, 1),
    ("", 6, False, BLACK, 0),
    ("Data Layer:  MySQL 8 (primary)  |  Redis 7 (cache, queues, sessions)", 15, True, CRIMSON, 0),
    ("", 6, False, BLACK, 0),
    ("External:  Payment Gateway  |  Google OAuth 2.0 + Calendar API  |  SMTP (configurable)", 15, True, CRIMSON, 0),
])
print("Architecture slide done")

# ── SLIDE: PSM 2 PLANNING ─────────────────────────────────────────────────────
slide = S["psm2"]
set_section_title(slide, "(8) PSM 2 Planning")
set_content(slide, [
    ("PSM II shifts focus entirely to system development, testing, and deployment.", 16, False, GREY, 0),
    ("", 8, False, BLACK, 0),
    ("Phase 1 — Environment Setup", 15, True, CRIMSON, 0),
    ("Configure Laragon local server, MySQL, Redis, and all Composer + npm dependencies", 14, False, DARK, 1),
    ("", 6, False, BLACK, 0),
    ("Phase 2 — Backend Development (Agile sprints)", 15, True, CRIMSON, 0),
    ("Implement all 11 modules in dependency order using Laravel 11", 14, False, DARK, 1),
    ("Integrate payment gateway (webhooks), Laravel Reverb (WebSocket), QR code generation", 14, False, DARK, 1),
    ("", 6, False, BLACK, 0),
    ("Phase 3 — Frontend Development (concurrent with backend)", 15, True, CRIMSON, 0),
    ("Implement wireframes as functional Blade + Alpine.js pages with Bootstrap 5", 14, False, DARK, 1),
    ("", 6, False, BLACK, 0),
    ("Phase 4 — Testing", 15, True, CRIMSON, 0),
    ("Unit tests, integration tests, and User Acceptance Testing (UAT)", 14, False, DARK, 1),
    ("", 6, False, BLACK, 0),
    ("Phase 5 — Deployment", 15, True, CRIMSON, 0),
    ("Deploy on Nginx + PHP-FPM + MySQL 8 + Redis; enable SSL with Let's Encrypt", 14, False, DARK, 1),
])
print("PSM2 Planning slide done")

# ── SLIDE: CONCLUSION ─────────────────────────────────────────────────────────
slide = S["conclusion"]
set_section_title(slide, "(9) Conclusion")
set_content(slide, [
    ("Problems that led to HomeLodge:", 16, True, CRIMSON, 0),
    ("Double-bookings, unstructured payment, physical key dependency, scattered communication across chat threads", 15, False, DARK, 1),
    ("", 8, False, BLACK, 0),
    ("Existing platforms are insufficient:", 16, True, CRIMSON, 0),
    ("Airbnb, Booking.com, and Agoda handle availability and payment but operate as public marketplaces.", 15, False, DARK, 1),
    ("None provides QR door access, booking extensions, fine-grained RBAC, or commission-free private operation.", 15, False, DARK, 1),
    ("", 8, False, BLACK, 0),
    ("PSM I deliverables completed:", 16, True, CRIMSON, 0),
    ("Complete requirements specification (URS + PRD)", 15, False, DARK, 1),
    ("Full system design: architecture, ER diagram (16 tables), UI wireframes", 15, False, DARK, 1),
    ("47 use cases across 12 modules documented and traced to requirements", 15, False, DARK, 1),
    ("Hybrid Waterfall + Agile methodology selected and justified", 15, False, DARK, 1),
    ("", 8, False, BLACK, 0),
    ("PSM II will implement, test, and deploy HomeLodge.", 16, True, CRIMSON, 0),
])
print("Conclusion slide done")

# ── SLIDE: THANK YOU ──────────────────────────────────────────────────────────
# Keep template content as-is; it already has "Thank You" and "Q&A"
print("Thank You slide kept as template")

# ── SAVE ──────────────────────────────────────────────────────────────────────
prs.save(OUTPUT)
print(f"\nSaved: {OUTPUT}")
print(f"Total slides: {len(prs.slides)}")
