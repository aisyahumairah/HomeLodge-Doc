from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ─── Colour palette ──────────────────────────────────────────────────────────
NAVY      = RGBColor(0x0D, 0x1B, 0x3E)   # deep navy – slide backgrounds
DARK_BLUE = RGBColor(0x1A, 0x35, 0x6E)   # header boxes
ACCENT    = RGBColor(0x3B, 0x82, 0xF6)   # blue accent
GOLD      = RGBColor(0xF5, 0xC5, 0x18)   # gold highlight
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG  = RGBColor(0xF0, 0xF4, 0xFF)   # light slide bg
DARK_TEXT = RGBColor(0x1A, 0x1A, 0x2E)
MID_BLUE  = RGBColor(0x22, 0x55, 0xBB)

# Slide dimensions: widescreen 16:9
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

BLANK = prs.slide_layouts[6]   # completely blank layout

# ─── Helper utilities ─────────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill_rgb=None, alpha=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)   # MSO_SHAPE.RECTANGLE
    shape.line.fill.background()
    if fill_rgb:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_rgb
    else:
        shape.fill.background()
    return shape

def add_textbox(slide, text, x, y, w, h,
                font_size=18, bold=False, color=WHITE,
                align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf    = txBox.text_frame
    tf.word_wrap = wrap
    p  = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(font_size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic    = italic
    return txBox

def add_paragraph(tf, text, font_size=16, bold=False, color=DARK_TEXT,
                  align=PP_ALIGN.LEFT, space_before=6, italic=False,
                  level=0):
    p = tf.add_paragraph()
    p.alignment   = align
    p.space_before = Pt(space_before)
    p.level = level
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(font_size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic    = italic
    return p

def navy_slide_base(slide):
    """Full navy background."""
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill_rgb=NAVY)

def header_bar(slide, title, y=Inches(0), h=Inches(1.1)):
    """Dark blue header band across the top."""
    add_rect(slide, 0, y, SLIDE_W, h, fill_rgb=DARK_BLUE)
    add_textbox(slide, title,
                Inches(0.4), y + Inches(0.15),
                Inches(12.5), Inches(0.85),
                font_size=32, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

def gold_accent_bar(slide, y, h=Inches(0.06)):
    add_rect(slide, 0, y, SLIDE_W, h, fill_rgb=GOLD)

def slide_number(slide, num, total=10):
    add_textbox(slide, f"{num} / {total}",
                Inches(11.8), Inches(7.1), Inches(1.2), Inches(0.35),
                font_size=11, bold=False, color=RGBColor(0xAA,0xBB,0xFF),
                align=PP_ALIGN.RIGHT)

def footer_brand(slide):
    add_textbox(slide, "HomeLodge  |  PSM I",
                Inches(0.3), Inches(7.1), Inches(4), Inches(0.35),
                font_size=11, color=RGBColor(0xAA,0xBB,0xFF))

def content_box(slide, x, y, w, h, fill=RGBColor(0x14,0x2A,0x5E), radius=None):
    box = add_rect(slide, x, y, w, h, fill_rgb=fill)
    return box

def table_row(slide, cols, x, y, row_h, widths, fill, font_size=13, bold=False, color=WHITE):
    cx = x
    for i, (text, w) in enumerate(zip(cols, widths)):
        add_rect(slide, cx, y, w, row_h, fill_rgb=fill)
        add_textbox(slide, text, cx + Inches(0.08), y + Pt(4), w - Inches(0.16), row_h,
                    font_size=font_size, bold=bold, color=color, wrap=True)
        cx += w

# ─── SLIDE 1: TITLE ──────────────────────────────────────────────────────────
def slide1(prs):
    slide = prs.slides.add_slide(BLANK)
    navy_slide_base(slide)

    # Large accent rectangle on right
    add_rect(slide, Inches(8.5), 0, Inches(4.83), SLIDE_H, fill_rgb=DARK_BLUE)
    # Gold vertical stripe
    add_rect(slide, Inches(8.35), 0, Inches(0.12), SLIDE_H, fill_rgb=GOLD)

    # Decorative circles
    for r, cx, cy, col in [
        (Inches(3), Inches(10.5), Inches(1.5), RGBColor(0x1A,0x35,0x6E)),
        (Inches(1.8), Inches(11.8), Inches(5.5), RGBColor(0x0A,0x14,0x30)),
    ]:
        c = slide.shapes.add_shape(9, cx - r/2, cy - r/2, r, r)  # oval
        c.fill.solid(); c.fill.fore_color.rgb = col
        c.line.fill.background()

    # Title text
    add_textbox(slide, "HomeLodge",
                Inches(0.5), Inches(1.8), Inches(7.6), Inches(1.2),
                font_size=52, bold=True, color=WHITE)
    add_textbox(slide, "A Web-Based Homestay Booking\nManagement System for Multi-Unit Operators",
                Inches(0.5), Inches(3.0), Inches(7.6), Inches(1.5),
                font_size=22, bold=False, color=RGBColor(0xBB,0xCC,0xFF))
    # Gold underline
    add_rect(slide, Inches(0.5), Inches(4.6), Inches(4.5), Inches(0.07), fill_rgb=GOLD)

    # Info block
    for i, (label, val) in enumerate([
        ("Student",    "[ Your Name ]"),
        ("Supervisor", "[ Supervisor Name ]"),
        ("Programme",  "Bachelor of Software Engineering"),
        ("Session",    "2025 / 2026"),
    ]):
        iy = Inches(5.0) + i * Inches(0.52)
        add_textbox(slide, label + ":", Inches(0.5), iy, Inches(1.6), Inches(0.5),
                    font_size=13, bold=True, color=GOLD)
        add_textbox(slide, val,  Inches(2.1), iy, Inches(6.0), Inches(0.5),
                    font_size=13, bold=False, color=WHITE)

    slide_number(slide, 1)
    return slide

# ─── SLIDE 2: PROJECT INTRODUCTION ───────────────────────────────────────────
def slide2(prs):
    slide = prs.slides.add_slide(BLANK)
    navy_slide_base(slide)
    header_bar(slide, "02  |  Project Introduction")
    gold_accent_bar(slide, Inches(1.1))

    subtitle = "When one homestay is fine. Three or more — that's where things break."
    add_textbox(slide, subtitle, Inches(0.4), Inches(1.2), Inches(12.5), Inches(0.45),
                font_size=15, italic=True, color=RGBColor(0xBB,0xCC,0xFF))

    problems = [
        ("No shared availability record",  "Two guests book the same dates"),
        ("Manual payment collection",      "No receipts, no consistent refund process"),
        ("Physical key handover",          "Operator must be present at every check-in"),
        ("Scattered chat threads",         "No link between conversation and booking"),
        ("Manual reminders",               "Missed notifications, lapsed bookings"),
    ]

    col_w = [Inches(4.2), Inches(7.6)]
    widths = col_w

    # Header row
    table_row(slide,
              ["Problem", "What Goes Wrong"],
              Inches(0.4), Inches(1.75),
              Inches(0.5), widths,
              fill=DARK_BLUE, font_size=14, bold=True, color=GOLD)

    fills = [RGBColor(0x14,0x2A,0x5E), RGBColor(0x0D,0x1B,0x3E)]
    for i, (prob, effect) in enumerate(problems):
        ry = Inches(1.75) + Inches(0.5) + i * Inches(0.78)
        table_row(slide, [prob, effect], Inches(0.4), ry,
                  Inches(0.75), widths,
                  fill=fills[i % 2], font_size=13, bold=False, color=WHITE)

    footer_brand(slide)
    slide_number(slide, 2)
    return slide

# ─── SLIDE 3: PROJECT AIM ─────────────────────────────────────────────────────
def slide3(prs):
    slide = prs.slides.add_slide(BLANK)
    navy_slide_base(slide)
    header_bar(slide, "03  |  Project Aim")
    gold_accent_bar(slide, Inches(1.1))

    # Central aim box
    content_box(slide, Inches(0.5), Inches(1.35), Inches(12.3), Inches(2.3),
                fill=RGBColor(0x14,0x2A,0x5E))
    add_textbox(slide, "AIM", Inches(0.7), Inches(1.45), Inches(2), Inches(0.5),
                font_size=13, bold=True, color=GOLD)
    add_textbox(slide,
        "To develop a web-based homestay booking management system that handles "
        "reservations, payment, property access, and guest communication — all in one "
        "place — for operators managing multiple homestay units.",
        Inches(0.7), Inches(1.85), Inches(12.0), Inches(1.7),
        font_size=20, bold=False, color=WHITE, wrap=True)

    # Key distinction
    add_rect(slide, Inches(0.5), Inches(3.85), Inches(12.3), Inches(0.08), fill_rgb=GOLD)
    add_textbox(slide, "Not another marketplace.",
                Inches(0.5), Inches(4.0), Inches(6), Inches(0.7),
                font_size=26, bold=True, color=GOLD)
    add_textbox(slide, "A private management tool that the operator owns, configures, and controls.",
                Inches(0.5), Inches(4.65), Inches(12.3), Inches(0.6),
                font_size=18, bold=False, color=RGBColor(0xBB,0xCC,0xFF))

    # 3 icons row
    icons = [
        ("No commission",        "Zero per-booking marketplace fees"),
        ("Full control",         "Operator owns the system and all data"),
        ("Private listing",      "No mandatory public marketplace exposure"),
    ]
    for i, (head, sub) in enumerate(icons):
        bx = Inches(0.5) + i * Inches(4.25)
        content_box(slide, bx, Inches(5.5), Inches(4.0), Inches(1.7),
                    fill=DARK_BLUE)
        add_textbox(slide, head, bx + Inches(0.15), Inches(5.6),
                    Inches(3.7), Inches(0.55), font_size=17, bold=True, color=GOLD)
        add_textbox(slide, sub, bx + Inches(0.15), Inches(6.1),
                    Inches(3.7), Inches(0.9), font_size=13, color=WHITE, wrap=True)

    footer_brand(slide)
    slide_number(slide, 3)
    return slide

# ─── SLIDE 4: PROJECT OBJECTIVES ──────────────────────────────────────────────
def slide4(prs):
    slide = prs.slides.add_slide(BLANK)
    navy_slide_base(slide)
    header_bar(slide, "04  |  Project Objectives")
    gold_accent_bar(slide, Inches(1.1))

    objectives = [
        ("01", "Study & Analyse",
         "Gather and document all functional and non-functional requirements for HomeLodge."),
        ("02", "Design",
         "Produce the system architecture, database schema, and UI covering all modules and both user roles."),
        ("03", "Develop",
         "Build the HomeLodge web application with 11 functional modules as defined by the requirements."),
        ("04", "Test",
         "Verify correctness, security, and usability against the stated requirements."),
    ]

    for i, (num, title, desc) in enumerate(objectives):
        bx = Inches(0.35) + (i % 2) * Inches(6.5)
        by = Inches(1.3) + (i // 2) * Inches(2.5)
        content_box(slide, bx, by, Inches(6.2), Inches(2.3), fill=RGBColor(0x14,0x2A,0x5E))
        # Number badge
        badge = slide.shapes.add_shape(9, bx + Inches(0.15), by + Inches(0.18),
                                       Inches(0.7), Inches(0.7))
        badge.fill.solid(); badge.fill.fore_color.rgb = GOLD
        badge.line.fill.background()
        add_textbox(slide, num, bx + Inches(0.18), by + Inches(0.22),
                    Inches(0.65), Inches(0.5), font_size=15, bold=True,
                    color=DARK_TEXT, align=PP_ALIGN.CENTER)
        add_textbox(slide, title, bx + Inches(1.0), by + Inches(0.2),
                    Inches(5.0), Inches(0.55), font_size=20, bold=True, color=GOLD)
        add_textbox(slide, desc, bx + Inches(0.15), by + Inches(0.85),
                    Inches(5.9), Inches(1.3), font_size=14, color=WHITE, wrap=True)

    # 11 modules strip
    add_rect(slide, 0, Inches(6.4), SLIDE_W, Inches(0.85), fill_rgb=DARK_BLUE)
    modules_text = ("11 Modules:  Authentication · Homestay Management · Booking · Payment · "
                    "Notification · Chat · User Management · Role & Permission · "
                    "System Settings · Audit Logs · QR Code Door Access")
    add_textbox(slide, modules_text, Inches(0.3), Inches(6.45),
                Inches(12.8), Inches(0.75), font_size=12, color=GOLD, wrap=True)

    footer_brand(slide)
    slide_number(slide, 4)
    return slide

# ─── SLIDE 5: BACKGROUND STUDY ────────────────────────────────────────────────
def slide5(prs):
    slide = prs.slides.add_slide(BLANK)
    navy_slide_base(slide)
    header_bar(slide, "05  |  Background Study — Why Existing Systems Fall Short")
    gold_accent_bar(slide, Inches(1.1))

    # Comparison table
    features = [
        ("Online availability calendar",     True,  True,  True,  True),
        ("Online payment processing",        True,  True,  True,  True),
        ("In-app guest-to-host messaging",   True,  True,  True,  True),
        ("QR code door access",              False, False, False, True),
        ("Booking extension management",     False, False, False, True),
        ("Role-based access control",        False, False, False, True),
        ("Private, non-marketplace",         False, False, False, True),
        ("Immutable audit log",              False, False, False, True),
        ("No per-booking commission",        False, False, False, True),
    ]

    headers = ["Feature", "Airbnb", "Booking.com", "Agoda", "HomeLodge"]
    col_ws  = [Inches(4.8), Inches(1.85), Inches(1.85), Inches(1.85), Inches(2.65)]

    table_row(slide, headers, Inches(0.3), Inches(1.25),
              Inches(0.55), col_ws, fill=DARK_BLUE, font_size=13, bold=True, color=GOLD)

    fills = [RGBColor(0x14,0x2A,0x5E), RGBColor(0x0D,0x1B,0x3E)]
    for i, row in enumerate(features):
        feat = row[0]
        vals = row[1:]
        ry = Inches(1.25) + Inches(0.55) + i * Inches(0.6)
        # Feature name cell
        add_rect(slide, Inches(0.3), ry, col_ws[0], Inches(0.58), fill_rgb=fills[i % 2])
        add_textbox(slide, feat, Inches(0.4), ry + Pt(2), col_ws[0] - Inches(0.15), Inches(0.56),
                    font_size=12, color=WHITE)
        # Platform cells
        cx = Inches(0.3) + col_ws[0]
        for j, (val, cw) in enumerate(zip(vals, col_ws[1:])):
            cell_fill = fills[i % 2]
            if j == 3:  # HomeLodge column
                cell_fill = RGBColor(0x0A,0x28,0x5A) if val else fills[i % 2]
            add_rect(slide, cx, ry, cw, Inches(0.58), fill_rgb=cell_fill)
            tick = "✓" if val else "✗"
            col  = RGBColor(0x4A,0xDE,0x80) if val else RGBColor(0xFF,0x55,0x55)
            if j == 3 and val:
                col = GOLD
            add_textbox(slide, tick, cx, ry + Pt(2), cw, Inches(0.56),
                        font_size=16, bold=True, color=col, align=PP_ALIGN.CENTER)
            cx += cw

    # Caption
    add_textbox(slide,
        "All three platforms solve the basics. None give operators a private, controllable system.",
        Inches(0.3), Inches(7.05), Inches(12.5), Inches(0.38),
        font_size=12, italic=True, color=RGBColor(0xBB,0xCC,0xFF))

    footer_brand(slide)
    slide_number(slide, 5)
    return slide

# ─── SLIDE 6: METHODOLOGY ────────────────────────────────────────────────────
def slide6(prs):
    slide = prs.slides.add_slide(BLANK)
    navy_slide_base(slide)
    header_bar(slide, "06  |  Development Methodology")
    gold_accent_bar(slide, Inches(1.1))

    phases = [
        ("01", "Requirements Analysis", "Waterfall",
         "Gathered, documented, and froze all requirements before writing any code."),
        ("02", "System Design",         "Waterfall",
         "Produced ER diagram, use case model, and UI wireframes — complete before coding."),
        ("03", "Implementation",        "Agile (Scrum + Kanban)",
         "11 modules built in time-boxed sprints; progress tracked per task on Kanban Board."),
        ("04", "Testing",               "Agile (Scrum + Kanban)",
         "Each module tested immediately after completion; defects fixed within the same sprint."),
        ("05", "Deployment",            "Waterfall",
         "Sequential checklist: Nginx, PHP-FPM, MySQL, Redis, SSL — all phases already done."),
    ]

    for i, (num, phase, method, desc) in enumerate(phases):
        by = Inches(1.3) + i * Inches(1.16)
        # connector line
        if i < len(phases) - 1:
            add_rect(slide, Inches(1.42), by + Inches(0.88), Inches(0.06), Inches(0.3),
                     fill_rgb=GOLD)

        # circle badge
        badge = slide.shapes.add_shape(9, Inches(1.0), by + Inches(0.05),
                                       Inches(0.82), Inches(0.82))
        badge.fill.solid(); badge.fill.fore_color.rgb = GOLD
        badge.line.fill.background()
        add_textbox(slide, num, Inches(1.0), by + Inches(0.12),
                    Inches(0.82), Inches(0.62), font_size=16, bold=True,
                    color=DARK_TEXT, align=PP_ALIGN.CENTER)

        # Phase box
        content_box(slide, Inches(2.05), by, Inches(11.0), Inches(0.95),
                    fill=RGBColor(0x14,0x2A,0x5E))
        add_textbox(slide, phase, Inches(2.2), by + Inches(0.08),
                    Inches(4.5), Inches(0.5), font_size=17, bold=True, color=WHITE)
        # Method badge
        mb = slide.shapes.add_shape(1, Inches(6.8), by + Inches(0.16),
                                    Inches(2.9), Inches(0.55))
        mb.fill.solid(); mb.fill.fore_color.rgb = DARK_BLUE
        mb.line.color.rgb = GOLD; mb.line.width = Pt(1)
        add_textbox(slide, method, Inches(6.85), by + Inches(0.18),
                    Inches(2.8), Inches(0.5), font_size=12, bold=True, color=GOLD)
        add_textbox(slide, desc, Inches(9.8), by + Inches(0.1),
                    Inches(3.15), Inches(0.8), font_size=11, color=WHITE, wrap=True)

    # Why box
    add_rect(slide, Inches(0.3), Inches(7.0), SLIDE_W - Inches(0.6), Inches(0.37),
             fill_rgb=DARK_BLUE)
    add_textbox(slide,
        "Why hybrid?  Waterfall → formal PSM documentation required.  "
        "Agile → 11 modules need incremental, trackable delivery.",
        Inches(0.4), Inches(7.02), Inches(12.5), Inches(0.35),
        font_size=12, color=GOLD)

    footer_brand(slide)
    slide_number(slide, 6)
    return slide

# ─── SLIDE 7: SYSTEM REQUIREMENTS ────────────────────────────────────────────
def slide7(prs):
    slide = prs.slides.add_slide(BLANK)
    navy_slide_base(slide)
    header_bar(slide, "07  |  System Requirements & Technology Stack")
    gold_accent_bar(slide, Inches(1.1))

    # Left: hardware
    add_textbox(slide, "Hardware (Development)", Inches(0.4), Inches(1.3), Inches(4), Inches(0.45),
                font_size=14, bold=True, color=GOLD)
    hw = [
        ("Laptop", "Intel Core i5, 8 GB RAM"),
        ("Storage", "256 GB SSD minimum"),
        ("Internet", "Broadband, min. 10 Mbps"),
        ("Server", "2 vCPU, 4 GB RAM, 20 GB SSD"),
    ]
    for i, (comp, spec) in enumerate(hw):
        by = Inches(1.75) + i * Inches(0.62)
        content_box(slide, Inches(0.4), by, Inches(5.8), Inches(0.57),
                    fill=RGBColor(0x14,0x2A,0x5E))
        add_textbox(slide, comp, Inches(0.55), by + Pt(4), Inches(1.5), Inches(0.5),
                    font_size=12, bold=True, color=GOLD)
        add_textbox(slide, spec, Inches(2.1), by + Pt(4), Inches(4.0), Inches(0.5),
                    font_size=12, color=WHITE)

    # Right: tech stack
    add_textbox(slide, "Technology Stack", Inches(6.8), Inches(1.3), Inches(6), Inches(0.45),
                font_size=14, bold=True, color=GOLD)
    stack = [
        ("Backend",     "Laravel 11  ·  PHP 8.2"),
        ("Frontend",    "Blade  ·  Alpine.js  ·  Bootstrap 5"),
        ("Database",    "MySQL 8  ·  Eloquent ORM"),
        ("Real-time",   "Laravel Reverb (WebSocket)"),
        ("Auth / RBAC", "laravel/socialite  ·  spatie/laravel-permission"),
        ("QR Code",     "simplesoftwareio/simple-qrcode (local)"),
        ("Payment",     "Billplz / Toyyibpay / Stripe (webhook)"),
    ]
    fills2 = [RGBColor(0x14,0x2A,0x5E), RGBColor(0x0D,0x1B,0x3E)]
    for i, (layer, tech) in enumerate(stack):
        by = Inches(1.75) + i * Inches(0.72)
        add_rect(slide, Inches(6.8), by, Inches(6.2), Inches(0.68),
                 fill_rgb=fills2[i % 2])
        add_textbox(slide, layer, Inches(6.95), by + Pt(4), Inches(1.8), Inches(0.6),
                    font_size=12, bold=True, color=GOLD)
        add_textbox(slide, tech, Inches(8.85), by + Pt(4), Inches(4.0), Inches(0.6),
                    font_size=12, color=WHITE)

    footer_brand(slide)
    slide_number(slide, 7)
    return slide

# ─── SLIDE 8: PSM I RESULTS ───────────────────────────────────────────────────
def slide8(prs):
    slide = prs.slides.add_slide(BLANK)
    navy_slide_base(slide)
    header_bar(slide, "08  |  PSM I Results — Analysis & Design Completed")
    gold_accent_bar(slide, Inches(1.1))

    # 4-quadrant layout
    quadrants = [
        ("Use Case Model",
         ["4 Actors: Guest, Admin, Payment Gateway, Google Calendar API",
          "12 Modules  ·  47 Use Cases documented",
          "All use cases traced to URS and PRD requirements"]),
        ("System Architecture",
         ["Client: Blade + Alpine.js + Bootstrap 5 + Laravel Echo",
          "Server: Nginx → PHP-FPM → Laravel 11 + Reverb",
          "Data: MySQL 8 (primary)  ·  Redis 7 (cache / queues)",
          "External: Payment Gateway, Google OAuth + Calendar, SMTP"]),
        ("Database Design",
         ["16 tables with foreign key constraints and soft deletes",
          "Key tables: bookings, booking_extensions, payments, qr_codes",
          "settings table: key-value model, no schema change needed",
          "Idempotent webhook handling via unique gateway_reference index"]),
        ("User Interface Design",
         ["Guest: top navbar, linear booking flow (browse → book → pay → QR)",
          "Admin: fixed left sidebar, 10 navigation groups",
          "Responsive: mobile / tablet / desktop breakpoints",
          "Design: Inter font, 8px grid, colour-coded booking status badges"]),
    ]

    positions = [
        (Inches(0.3),  Inches(1.25)),
        (Inches(6.75), Inches(1.25)),
        (Inches(0.3),  Inches(4.35)),
        (Inches(6.75), Inches(4.35)),
    ]
    box_w = Inches(6.2)
    box_h = Inches(2.85)

    for (heading, bullets), (bx, by) in zip(quadrants, positions):
        content_box(slide, bx, by, box_w, box_h, fill=RGBColor(0x14,0x2A,0x5E))
        add_textbox(slide, heading, bx + Inches(0.15), by + Inches(0.1),
                    box_w - Inches(0.3), Inches(0.5),
                    font_size=16, bold=True, color=GOLD)
        add_rect(slide, bx + Inches(0.15), by + Inches(0.6),
                 box_w - Inches(0.3), Pt(1.5), fill_rgb=GOLD)
        for j, bullet in enumerate(bullets):
            add_textbox(slide, "  •  " + bullet,
                        bx + Inches(0.15), by + Inches(0.7) + j * Inches(0.53),
                        box_w - Inches(0.3), Inches(0.55),
                        font_size=12, color=WHITE, wrap=True)

    footer_brand(slide)
    slide_number(slide, 8)
    return slide

# ─── SLIDE 9: CONCLUSION ──────────────────────────────────────────────────────
def slide9(prs):
    slide = prs.slides.add_slide(BLANK)
    navy_slide_base(slide)
    header_bar(slide, "09  |  Conclusion")
    gold_accent_bar(slide, Inches(1.1))

    # Problems identified
    add_textbox(slide, "Problems Identified", Inches(0.4), Inches(1.3), Inches(5.8), Inches(0.45),
                font_size=14, bold=True, color=GOLD)
    probs = [
        "Double-bookings from manually tracking availability",
        "Unstructured payment with no consistent refund process",
        "Physical key handover requiring operator presence at every check-in",
        "Communication scattered across personal chats with no booking link",
    ]
    for i, p in enumerate(probs):
        by = Inches(1.75) + i * Inches(0.72)
        content_box(slide, Inches(0.4), by, Inches(5.8), Inches(0.65),
                    fill=RGBColor(0x14,0x2A,0x5E))
        add_textbox(slide, "✗  " + p, Inches(0.55), by + Pt(5),
                    Inches(5.55), Inches(0.6), font_size=12.5,
                    color=WHITE, wrap=True)

    # Vertical divider
    add_rect(slide, Inches(6.45), Inches(1.25), Inches(0.06), Inches(5.7), fill_rgb=GOLD)

    # PSM I achievements
    add_textbox(slide, "PSM I Delivered", Inches(6.7), Inches(1.3), Inches(6), Inches(0.45),
                font_size=14, bold=True, color=GOLD)
    psm1 = [
        "Complete requirements spec (URS + PRD)",
        "Full system design: architecture, ER diagram, UI wireframes",
        "47 use cases across 12 modules documented and traced",
        "Hybrid methodology selected and justified (Waterfall + Agile)",
        "Technology stack selected and compared",
    ]
    for i, item in enumerate(psm1):
        by = Inches(1.75) + i * Inches(0.68)
        content_box(slide, Inches(6.7), by, Inches(6.3), Inches(0.62),
                    fill=RGBColor(0x14,0x2A,0x5E))
        add_textbox(slide, "✓  " + item, Inches(6.85), by + Pt(5),
                    Inches(6.05), Inches(0.58), font_size=12.5,
                    color=RGBColor(0x4A,0xDE,0x80), wrap=True)

    # PSM II plan strip
    add_rect(slide, 0, Inches(6.45), SLIDE_W, Inches(0.85), fill_rgb=DARK_BLUE)
    add_textbox(slide, "PSM II Plan:", Inches(0.3), Inches(6.5), Inches(1.5), Inches(0.7),
                font_size=13, bold=True, color=GOLD)
    add_textbox(slide,
        "Implement 11 modules  ·  Integrate payment gateway, WebSocket chat, QR door access  "
        "·  Unit, integration & UAT testing  ·  Deploy on Nginx + MySQL + Redis",
        Inches(1.9), Inches(6.5), Inches(11.0), Inches(0.7),
        font_size=13, color=WHITE, wrap=True)

    footer_brand(slide)
    slide_number(slide, 9)
    return slide

# ─── SLIDE 10: THANK YOU ──────────────────────────────────────────────────────
def slide10(prs):
    slide = prs.slides.add_slide(BLANK)
    navy_slide_base(slide)

    # Large dark blue right panel
    add_rect(slide, Inches(8.0), 0, Inches(5.33), SLIDE_H, fill_rgb=DARK_BLUE)
    add_rect(slide, Inches(7.85), 0, Inches(0.12), SLIDE_H, fill_rgb=GOLD)

    # Decorative element
    c = slide.shapes.add_shape(9, Inches(9.5), Inches(0.5), Inches(3.5), Inches(3.5))
    c.fill.solid(); c.fill.fore_color.rgb = RGBColor(0x0A,0x14,0x30)
    c.line.fill.background()

    add_textbox(slide, "Thank You",
                Inches(0.5), Inches(1.5), Inches(7.0), Inches(1.2),
                font_size=60, bold=True, color=WHITE)
    add_rect(slide, Inches(0.5), Inches(2.8), Inches(5.5), Inches(0.08), fill_rgb=GOLD)
    add_textbox(slide, "Open for Questions",
                Inches(0.5), Inches(3.0), Inches(7.0), Inches(0.7),
                font_size=26, bold=False, color=RGBColor(0xBB,0xCC,0xFF))

    add_textbox(slide, "HomeLodge",
                Inches(0.5), Inches(4.0), Inches(7.0), Inches(0.75),
                font_size=28, bold=True, color=GOLD)
    add_textbox(slide, "A Web-Based Homestay Booking Management System",
                Inches(0.5), Inches(4.65), Inches(7.0), Inches(0.55),
                font_size=16, color=WHITE)

    for i, (label, val) in enumerate([
        ("Student",    "[ Your Name ]"),
        ("Supervisor", "[ Supervisor Name ]"),
        ("Session",    "2025 / 2026"),
    ]):
        iy = Inches(5.45) + i * Inches(0.55)
        add_textbox(slide, label + ":", Inches(0.5), iy, Inches(1.5), Inches(0.5),
                    font_size=14, bold=True, color=GOLD)
        add_textbox(slide, val, Inches(2.0), iy, Inches(5.5), Inches(0.5),
                    font_size=14, color=WHITE)

    # Right panel: PSM summary stats
    stats = [
        ("47",  "Use Cases"),
        ("12",  "Modules"),
        ("16",  "DB Tables"),
        ("11",  "Functional\nModules"),
    ]
    for i, (num, label) in enumerate(stats):
        bx = Inches(8.15) + (i % 2) * Inches(2.55)
        by = Inches(1.8) + (i // 2) * Inches(2.5)
        content_box(slide, bx, by, Inches(2.3), Inches(2.1),
                    fill=RGBColor(0x0D,0x1B,0x3E))
        add_textbox(slide, num, bx + Inches(0.1), by + Inches(0.2),
                    Inches(2.1), Inches(1.1), font_size=52, bold=True,
                    color=GOLD, align=PP_ALIGN.CENTER)
        add_textbox(slide, label, bx + Inches(0.1), by + Inches(1.35),
                    Inches(2.1), Inches(0.65), font_size=13, color=WHITE,
                    align=PP_ALIGN.CENTER, wrap=True)

    slide_number(slide, 10)
    return slide

# ─── BUILD PRESENTATION ───────────────────────────────────────────────────────
slide1(prs)
slide2(prs)
slide3(prs)
slide4(prs)
slide5(prs)
slide6(prs)
slide7(prs)
slide8(prs)
slide9(prs)
slide10(prs)

out = r"c:\Users\Admin\Documents\HomeLodge-Doc\docs\Presentation\HomeLodge_PSM1_Presentation.pptx"
prs.save(out)
print(f"Saved: {out}")
